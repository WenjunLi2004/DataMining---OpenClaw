"""Build a presentation deck for the OpenClaw open-source radar project.

The deck is optimized for live course presentation: one claim per slide,
large text, and visual proof objects instead of report-style paragraphs.
"""
import os
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


OUT = os.environ.get("PPT_OUT", "/Users/wenjun/openclaw-project/reports/temp_presentation.pptx")
TOTAL = 10
FONT = "LXGW WenKai"

INK = RGBColor(0x24, 0x29, 0x2F)
MUTED = RGBColor(0x57, 0x60, 0x6A)
BORDER = RGBColor(0xD0, 0xD7, 0xDE)
BG = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE = RGBColor(0xF6, 0xF8, 0xFA)
GREEN = RGBColor(0x2D, 0xA4, 0x4E)
BLUE = RGBColor(0x09, 0x69, 0xDA)
ORANGE = RGBColor(0xFB, 0x85, 0x00)
PURPLE = RGBColor(0x8A, 0x3F, 0xFC)
DARK = RGBColor(0x1F, 0x23, 0x2A)
RED = RGBColor(0xCF, 0x22, 0x2E)


def emu(value):
    return Emu(int(value))


def bg(slide, color=BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def text(slide, x, y, w, h, body, *, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = body
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def multi(slide, x, y, w, h, lines, *, line_spacing=1.12, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    for i, (body, opts) in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = body
        run.font.name = opts.get("font", FONT)
        run.font.size = Pt(opts.get("size", 18))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", INK)
    return box


def rect(slide, x, y, w, h, *, fill=None, line=BORDER, line_w=0.8, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        s.adjustments[0] = 0.10
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def header(slide, title, idx):
    text(slide, Inches(0.62), Inches(0.28), Inches(10.6), Inches(0.45),
         title, size=25, bold=True)
    text(slide, Inches(11.55), Inches(0.34), Inches(1.1), Inches(0.3),
         f"{idx} / {TOTAL}", size=12, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)
    rect(slide, Inches(0.62), Inches(0.92), Inches(12.1), Inches(0.012),
         fill=BORDER, line=None)


def footer(slide, label):
    text(slide, Inches(0.62), Inches(7.12), Inches(7), Inches(0.22),
         "OpenClaw · GitHub Repository Potential Radar", size=9, color=MUTED)
    text(slide, Inches(9.2), Inches(7.12), Inches(3.5), Inches(0.22),
         label, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def stat(slide, x, y, w, value, label, color):
    rect(slide, x, y, w, Inches(0.86), fill=SUBTLE, line=BORDER, radius=True)
    text(slide, x, y + Inches(0.10), w, Inches(0.34), value,
         size=22, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT)
    text(slide, x, y + Inches(0.52), w, Inches(0.24), label,
         size=10, color=MUTED, align=PP_ALIGN.CENTER, font=FONT)


def step(slide, x, y, n, title, desc, color):
    rect(slide, x, y, Inches(1.00), Inches(0.64), fill=color, line=None, radius=True)
    text(slide, x, y + Inches(0.14), Inches(1.00), Inches(0.32), n,
         size=16, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
    text(slide, x + Inches(1.22), y - Inches(0.02), Inches(3.8), Inches(0.34),
         title, size=17, bold=True)
    text(slide, x + Inches(1.22), y + Inches(0.38), Inches(4.0), Inches(0.32),
         desc, size=13, color=MUTED)


def bar(slide, x, y, label, value, max_value, color, *, value_label=None, width=3.1):
    text(slide, x, y, Inches(1.15), Inches(0.28), label, size=13, color=INK, font=FONT)
    rect(slide, x + Inches(1.28), y + Inches(0.06), Inches(width), Inches(0.18),
         fill=SUBTLE, line=None)
    rect(slide, x + Inches(1.28), y + Inches(0.06),
         emu(Inches(width) * (value / max_value)), Inches(0.18),
         fill=color, line=None)
    text(slide, x + Inches(1.28 + width + 0.10), y - Inches(0.01), Inches(0.65), Inches(0.28),
         value_label or str(value), size=12, bold=True, color=INK, font=FONT)


def pill(slide, x, y, w, label, color):
    rect(slide, x, y, w, Inches(0.36), fill=SUBTLE, line=BORDER, radius=True)
    text(slide, x, y + Inches(0.075), w, Inches(0.18), label,
         size=10, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT)


def star(slide, x, y, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, x, y, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def circle(slide, x, y, size, *, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.8)
    s.shadow.inherit = False
    return s


def cover_node(slide, x, y, w, icon, title, sub, color):
    rect(slide, x, y, w, Inches(0.96), fill=BG, line=BORDER, radius=True)
    circle(slide, x + Inches(0.16), y + Inches(0.19), Inches(0.44), fill=color)
    text(slide, x + Inches(0.16), y + Inches(0.275), Inches(0.44), Inches(0.18),
         icon, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
    text(slide, x + Inches(0.72), y + Inches(0.18), w - Inches(0.86), Inches(0.26),
         title, size=15, bold=True)
    text(slide, x + Inches(0.72), y + Inches(0.55), w - Inches(0.86), Inches(0.20),
         sub, size=9.5, color=MUTED)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# Slide 1: Cover
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, Inches(0.18), Inches(7.5), fill=INK, line=None)
logo = "/Users/wenjun/openclaw-project/dashboard/lobster_logo.png"
github_logo = "/Users/wenjun/.codex/plugins/cache/openai-curated/github/7955f1db/assets/github.png"
s.shapes.add_picture(logo, Inches(10.44), Inches(0.78), width=Inches(1.38), height=Inches(1.38))
s.shapes.add_picture(github_logo, Inches(11.62), Inches(0.68), width=Inches(1.28), height=Inches(1.28))
star(s, Inches(12.40), Inches(2.04), Inches(0.22), ORANGE)
star(s, Inches(10.08), Inches(2.24), Inches(0.16), PURPLE)
star(s, Inches(11.84), Inches(2.30), Inches(0.13), GREEN)
text(s, Inches(0.78), Inches(0.70), Inches(6.3), Inches(0.34),
     "OpenClaw Project · Data Mining", size=13, color=MUTED)
text(s, Inches(0.78), Inches(1.48), Inches(8.9), Inches(0.76),
     "开源项目潜力发现系统", size=42, bold=True)
text(s, Inches(0.78), Inches(2.32), Inches(9.6), Inches(0.48),
     "从 GitHub 早期信号中挖掘值得观察的开源项目", size=25, bold=True, color=INK)
rect(s, Inches(0.80), Inches(2.98), Inches(7.75), Inches(0.07), fill=GREEN, line=None)

flow = [
    ("01", "历史样本", "500 repos", GREEN),
    ("30d", "早期信号", "创建后 30 天", BLUE),
    ("RF", "模型排序", "AUC 0.878", ORANGE),
    ("OK", "事实洞察", "诊断 + 解释", PURPLE),
    ("★", "今日雷达", "近期候选", GREEN),
]
for i, (icon, title, sub, c) in enumerate(flow):
    x = Inches(0.78 + i * 2.35)
    cover_node(s, x, Inches(3.58), Inches(1.95), icon, title, sub, c)
    if i < len(flow) - 1:
        text(s, x + Inches(2.00), Inches(3.90), Inches(0.30), Inches(0.18),
             "→", size=16, bold=True, color=BORDER, align=PP_ALIGN.CENTER)
chips = [("500", "repos", GREEN), ("38", "features", BLUE), ("0.878", "RF AUC", ORANGE), ("0.90", "P@10", PURPLE)]
for i, (v, l, c) in enumerate(chips):
    stat(s, Inches(0.78 + i * 2.55), Inches(5.35), Inches(2.15), v, l, c)
star(s, Inches(10.72), Inches(5.44), Inches(0.18), ORANGE)
star(s, Inches(10.94), Inches(5.32), Inches(0.12), ORANGE)
text(s, Inches(0.78), Inches(6.72), Inches(5.2), Inches(0.26),
     "李文俊 2023150001", size=17, bold=True, color=INK)
text(s, Inches(5.25), Inches(6.76), Inches(3.2), Inches(0.24),
     "OpenClaw Project · Data Mining", size=15, color=MUTED, align=PP_ALIGN.CENTER)
text(s, Inches(9.55), Inches(6.76), Inches(2.75), Inches(0.24),
     date.today().strftime("%Y · %m · %d"), size=12, color=MUTED, align=PP_ALIGN.RIGHT)


# Slide 2: Motivation
s = prs.slides.add_slide(BLANK)
bg(s)
text(s, Inches(0.62), Inches(0.18), Inches(3.8), Inches(0.48),
     "Motivation", size=31, bold=True)
text(s, Inches(11.55), Inches(0.28), Inches(1.1), Inches(0.3),
     f"2 / {TOTAL}", size=12, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)
rect(s, Inches(0.62), Inches(0.88), Inches(12.1), Inches(0.012),
     fill=BORDER, line=None)

text(s, Inches(1.30), Inches(1.35), Inches(10.9), Inches(0.36),
     "目标更早一步 —— 在仓库创建后的早期窗口，识别值得关注的潜力项目", size=23,
     bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(0.72), Inches(2.15), Inches(4.30), Inches(1.58),
     fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(1.02), Inches(2.42), Inches(3.65), Inches(0.35),
     "↘ GitHub Trending", size=24, bold=True, color=ORANGE)
multi(s, Inches(1.02), Inches(3.02), Inches(3.75), Inches(0.52), [
    ("看见的是：已经获得关注的项目", {"size": 16, "bold": True, "color": INK}),
    ("问题是：可能错过早期机会", {"size": 15, "color": MUTED}),
], line_spacing=1.25)

text(s, Inches(5.70), Inches(2.62), Inches(1.1), Inches(0.30),
     "启发", size=21, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.20), Inches(3.05), Inches(2.80), Inches(0.16))
arrow.fill.solid()
arrow.fill.fore_color.rgb = BLUE
arrow.line.fill.background()
arrow.shadow.inherit = False

rect(s, Inches(8.30), Inches(2.15), Inches(4.30), Inches(1.58),
     fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(8.60), Inches(2.42), Inches(3.55), Inches(0.35),
     "✧ OpenClaw Radar", size=24, bold=True, color=GREEN)
multi(s, Inches(8.60), Inches(3.02), Inches(3.75), Inches(0.52), [
    ("观察的是：前 30 天的项目行为", {"size": 16, "bold": True, "color": INK}),
    ("输出的是：可追踪的候选清单", {"size": 15, "color": MUTED}),
], line_spacing=1.25)

text(s, Inches(0.72), Inches(4.22), Inches(7.8), Inches(0.34),
     "研究支撑：相关研究表明这条路是可行的", size=20, bold=True)

papers = [
    (BLUE, "Predicting the Popularity of\nGitHub Repositories", "PROMISE'16",
     "用 stars 时间序列 + 多元回归即可\n预测仓库流行度", "→ 流行度是可建模的量"),
    (ORANGE, "Understanding the Factors that\nImpact the Popularity", "ICSME'16",
     "2,279 仓库 4 类增长模式；\n语言、领域显著影响 stars", "→ 项目特征带预测信号"),
    (GREEN, "What Makes a Popular\nAcademic AI Repository?", "EMSE / ICSE'21",
     "1,149 学术 AI 仓库；21 特征中 11 个\n在热门 vs 冷门间显著差异", "→ 早期工程实践可识别"),
]
for i, (c, title, venue, body, takeaway) in enumerate(papers):
    x = Inches(0.72 + i * 4.18)
    rect(s, x, Inches(4.74), Inches(3.95), Inches(1.98), fill=BG, line=BORDER, radius=True)
    rect(s, x, Inches(4.74), Inches(3.95), Inches(0.06), fill=c, line=None)
    text(s, x + Inches(0.30), Inches(4.98), Inches(3.35), Inches(0.42),
         title, size=13.5, bold=True)
    text(s, x + Inches(0.30), Inches(5.62), Inches(3.35), Inches(0.25),
         venue, size=13, bold=True, italic=True, color=c)
    text(s, x + Inches(0.30), Inches(5.98), Inches(3.35), Inches(0.42),
         body, size=12.5, color=INK)
    text(s, x + Inches(0.30), Inches(6.45), Inches(3.35), Inches(0.22),
         takeaway, size=12.5, bold=True, color=c)
footer(s, "Motivation")


# Slide 3: Flow
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "系统总流程：先验证，再应用", 3)
text(s, Inches(0.72), Inches(1.18), Inches(10.9), Inches(0.46),
     "整个项目分两段：历史回测证明方法有效，Today Radar 把模型用于近期仓库。", size=22, bold=True)

nodes = [
    ("1", "数据采集", "GitHub snapshot", GREEN),
    ("2", "特征工程", "19 features", BLUE),
    ("3", "模型训练", "LR / RF / XGB", ORANGE),
    ("4", "实验评估", "AUC / P@10", PURPLE),
    ("5", "事实诊断", "diagnostic JSON", GREEN),
    ("6", "今日雷达", "candidate list", BLUE),
]
x0 = Inches(0.72)
y0 = Inches(2.32)
w = Inches(1.78)
gap = Inches(0.28)
for i, (n, name, sub, c) in enumerate(nodes):
    x = x0 + i * (w + gap)
    rect(s, x, y0, w, Inches(2.00), fill=SUBTLE, line=BORDER, radius=True)
    rect(s, x + Inches(0.18), y0 + Inches(0.22), Inches(0.46), Inches(0.46), fill=c, line=None, radius=True)
    text(s, x + Inches(0.18), y0 + Inches(0.30), Inches(0.46), Inches(0.20), n,
         size=11, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
    text(s, x + Inches(0.20), y0 + Inches(0.88), w - Inches(0.40), Inches(0.30), name,
         size=18, bold=True, align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.20), y0 + Inches(1.28), w - Inches(0.40), Inches(0.28), sub,
         size=11, color=MUTED, align=PP_ALIGN.CENTER, font=FONT)
    if i < len(nodes) - 1:
        rect(s, x + w + Inches(0.03), y0 + Inches(0.98), gap - Inches(0.06), Inches(0.03), fill=BORDER, line=None)

rect(s, Inches(0.72), Inches(5.18), Inches(5.55), Inches(0.90), fill=RGBColor(0xEE, 0xF7, 0xEE), line=GREEN, radius=True)
text(s, Inches(1.02), Inches(5.38), Inches(5.0), Inches(0.28),
     "Backtest：验证早期信号是否真的有用", size=17, bold=True, color=GREEN)
rect(s, Inches(6.75), Inches(5.18), Inches(5.55), Inches(0.90), fill=RGBColor(0xEE, 0xF4, 0xFF), line=BLUE, radius=True)
text(s, Inches(7.05), Inches(5.38), Inches(5.0), Inches(0.28),
     "Radar：把模型用于今天的候选发现", size=17, bold=True, color=BLUE)
footer(s, "System Flow")


# Slide 4: Data collection
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "数据采集：只看前 30 天，标签放在未来", 4)
text(s, Inches(0.72), Inches(1.12), Inches(10.9), Inches(0.42),
     "每条样本被拆成 snapshot 和 label；snapshot 字段（含 README 和 contributors）都按 created_at + 30d 严格回溯。", size=21, bold=True)

rect(s, Inches(0.72), Inches(1.86), Inches(3.65), Inches(4.35), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(1.02), Inches(2.10), Inches(3.0), Inches(0.34), "500 个仓库", size=26, bold=True)
text(s, Inches(1.02), Inches(2.56), Inches(3.0), Inches(0.28), "GitHub Search API · AI/CS 关键词", size=12, color=MUTED)
lang = [("Rust", 151, GREEN), ("Go", 118, GREEN), ("Python", 100, GREEN), ("JS", 100, GREEN), ("TS", 31, ORANGE)]
for i, (name, val, c) in enumerate(lang):
    bar(s, Inches(1.02), Inches(3.05 + i * 0.42), name, val, 151, c, width=1.75)

rect(s, Inches(4.70), Inches(1.86), Inches(3.80), Inches(4.35), fill=BG, line=BORDER, radius=True)
text(s, Inches(5.00), Inches(2.10), Inches(3.2), Inches(0.34), "Snapshot 特征", size=22, bold=True, color=BLUE)
multi(s, Inches(5.00), Inches(2.72), Inches(3.05), Inches(2.3), [
    ("README (30d) 长度 / 图片 / Demo", {"size": 16}),
    ("commits / issues / PRs (30d)", {"size": 16}),
    ("contributors_30d (commit-author 去重)", {"size": 14}),
    ("语言 / owner 类型", {"size": 16}),
], line_spacing=1.35)
pill(s, Inches(5.0), Inches(5.35), Inches(2.80), "created_at → T+30d", BLUE)

rect(s, Inches(8.82), Inches(1.86), Inches(3.80), Inches(4.35), fill=BG, line=BORDER, radius=True)
text(s, Inches(9.12), Inches(2.10), Inches(3.2), Inches(0.34), "Label", size=22, bold=True, color=ORANGE)
multi(s, Inches(9.12), Inches(2.76), Inches(3.0), Inches(1.7), [
    ("current_stars", {"size": 20, "bold": True, "font": FONT}),
    ("is_top20 = stars ≥ p80", {"size": 18, "bold": True}),
    ("本批次 p80 = 48 stars", {"size": 16, "color": MUTED}),
], line_spacing=1.28)
pill(s, Inches(9.12), Inches(5.35), Inches(2.80), "训练时不进入 snapshot", ORANGE)
footer(s, "Data")


# Slide 5: Feature engineering
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "特征工程：把早期仓库行为变成 19 维向量（strict-30d v3）", 5)
text(s, Inches(0.72), Inches(1.12), Inches(10.9), Inches(0.42),
     "所有 19 个特征都能严格回溯到 created_at + 30d 内；当前态字段已主动移出模型。", size=21, bold=True)

groups = [
    ("6", "语言 one-hot", "Python · JS · Go · Rust · TS · Other", BLUE),
    ("1", "Owner 类型", "is_org", PURPLE),
    ("4", "早期活动 (30d)", "commits · issues · PRs · contributors", ORANGE),
    ("4", "历史 README (30d)", "len · has · image · demo URL", GREEN),
    ("4", "派生活跃度", "total · per_contrib · pr/issue · has_PR", INK),
]
for i, (num, title, desc, c) in enumerate(groups):
    x = Inches(0.72 + (i % 3) * 4.1)
    y = Inches(2.00 + (i // 3) * 1.55)
    width = Inches(3.68)
    rect(s, x, y, width, Inches(1.18), fill=SUBTLE, line=BORDER, radius=True)
    text(s, x + Inches(0.25), y + Inches(0.22), Inches(0.75), Inches(0.45),
         num, size=28, bold=True, color=c, font=FONT)
    text(s, x + Inches(1.10), y + Inches(0.22), width - Inches(1.3), Inches(0.30),
         title, size=17, bold=True)
    text(s, x + Inches(1.10), y + Inches(0.62), width - Inches(1.3), Inches(0.26),
         desc, size=11.5, color=MUTED)

rect(s, Inches(0.72), Inches(5.45), Inches(11.9), Inches(0.82), fill=DARK, line=None, radius=True)
text(s, Inches(1.02), Inches(5.65), Inches(2.2), Inches(0.32), "19 features", size=22, bold=True, color=BG, font=FONT)
text(s, Inches(3.25), Inches(5.69), Inches(8.8), Inches(0.26),
     "全部来自仓库创建后前 30 天，训练和 Today Radar 使用同一套 feature schema", size=15, color=BG)
footer(s, "Features")


# Slide 6: Model training
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "模型训练：三类模型对比，RF 作为候选排序器", 6)
text(s, Inches(0.72), Inches(1.12), Inches(11.3), Inches(0.42),
     "训练阶段不只是跑分，还要保存 artifacts，让模型能被 Today Radar 复用。", size=21, bold=True)

step(s, Inches(0.95), Inches(2.00), "01", "输入特征矩阵", "features.csv · 500 × 19 (strict-30d)", GREEN)
step(s, Inches(0.95), Inches(3.05), "02", "训练三类模型", "Logistic Regression · Random Forest · XGBoost", BLUE)
step(s, Inches(0.95), Inches(4.10), "03", "评估排序能力", "5-fold CV · AUC · PR-AUC · P@10", ORANGE)
step(s, Inches(0.95), Inches(5.15), "04", "保存复用产物", "rf_model.joblib · schema · thresholds", PURPLE)

rect(s, Inches(7.05), Inches(2.02), Inches(4.95), Inches(3.92), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(7.40), Inches(2.34), Inches(4.2), Inches(0.35), "为什么选 RF？", size=24, bold=True)
multi(s, Inches(7.40), Inches(3.02), Inches(4.25), Inches(1.65), [
    ("AUC 最高：0.878", {"size": 21, "bold": True, "color": GREEN, "font": FONT}),
    ("P@10 最高：0.90", {"size": 21, "bold": True, "color": GREEN, "font": FONT}),
    ("更适合“少量高质量候选推荐”", {"size": 17, "color": INK}),
], line_spacing=1.28)
text(s, Inches(7.40), Inches(5.15), Inches(4.2), Inches(0.34),
     "注意：RF recall 较低，所以它是 radar，不是全量发现器。", size=13, color=MUTED)
footer(s, "Training")


# Slide 7: Results
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "实验结果：早期信号确实能排序潜力仓库", 7)
text(s, Inches(0.72), Inches(1.12), Inches(11.1), Inches(0.42),
     "模型最有价值的不是“预测每一个项目”，而是把候选项目排到更值得看的顺序。", size=21, bold=True)

rect(s, Inches(0.72), Inches(1.92), Inches(3.55), Inches(1.55), fill=DARK, line=None, radius=True)
text(s, Inches(0.72), Inches(2.18), Inches(3.55), Inches(0.45), "0.878", size=36, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
text(s, Inches(0.72), Inches(2.84), Inches(3.55), Inches(0.28), "Random Forest AUC", size=13, color=GREEN, align=PP_ALIGN.CENTER, font=FONT)

rect(s, Inches(4.82), Inches(1.92), Inches(3.55), Inches(1.55), fill=DARK, line=None, radius=True)
text(s, Inches(4.82), Inches(2.18), Inches(3.55), Inches(0.45), "0.90", size=36, bold=True, color=BG, align=PP_ALIGN.CENTER, font=FONT)
text(s, Inches(4.82), Inches(2.84), Inches(3.55), Inches(0.28), "Precision@10", size=13, color=GREEN, align=PP_ALIGN.CENTER, font=FONT)

rect(s, Inches(8.92), Inches(1.92), Inches(3.55), Inches(1.55), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(8.92), Inches(2.22), Inches(3.55), Inches(0.36), "+0.066", size=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font=FONT)
text(s, Inches(8.92), Inches(2.86), Inches(3.55), Inches(0.28), "加入早期活动后的 AUC 提升", size=12, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, Inches(0.72), Inches(4.10), Inches(11.9), Inches(1.72), fill=BG, line=BORDER, radius=True)
text(s, Inches(1.02), Inches(4.35), Inches(2.1), Inches(0.30), "Ablation", size=15, bold=True, color=MUTED, font=FONT)
abl = [("Lang+Owner", 0.7855, MUTED), ("+ Activity 30d", 0.8516, GREEN), ("+ README 30d", 0.8756, BLUE), ("+ Derived", 0.8783, INK)]
base_x = Inches(2.45)
for i, (name, auc, c) in enumerate(abl):
    x = base_x + Inches(i * 2.38)
    text(s, x, Inches(4.32), Inches(1.8), Inches(0.25), name, size=13, bold=True)
    rect(s, x, Inches(4.82), Inches(1.7), Inches(0.22), fill=SUBTLE, line=None)
    rect(s, x, Inches(4.82), emu(Inches(1.7) * ((auc - 0.75) / 0.14)), Inches(0.22), fill=c, line=None)
    text(s, x, Inches(5.22), Inches(1.7), Inches(0.25), f"{auc:.4f}", size=14, bold=True, font=FONT)
text(s, Inches(1.02), Inches(6.25), Inches(11.2), Inches(0.32),
     "结论：30 天活跃度与历史 README 是主要信号；派生比率特征的边际增益较小。", size=19, bold=True)
footer(s, "Results")


# Slide 8: OpenClaw value
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "OpenClaw 的价值：不是更聪明，而是让复杂分析可复现", 8)
text(s, Inches(0.72), Inches(1.12), Inches(11.5), Inches(0.42),
     "单次问大模型也能写分析；OpenClaw 解决的是“多步骤任务如何稳定运行、留下证据、可重复检查”。", size=21, bold=True)

rect(s, Inches(0.78), Inches(2.02), Inches(5.55), Inches(3.70), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(1.12), Inches(2.34), Inches(4.8), Inches(0.34), "直接问大模型", size=24, bold=True, color=ORANGE)
multi(s, Inches(1.12), Inches(3.02), Inches(4.65), Inches(1.8), [
    ("输入结果文件，让模型总结", {"size": 18}),
    ("容易混入不存在的数字", {"size": 18, "color": RED}),
    ("每次回答不一定一致", {"size": 18, "color": RED}),
    ("难以接入完整 pipeline", {"size": 18, "color": RED}),
], line_spacing=1.28)

rect(s, Inches(7.00), Inches(2.02), Inches(5.55), Inches(3.70), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(7.34), Inches(2.34), Inches(4.8), Inches(0.34), "OpenClaw 工作流", size=24, bold=True, color=GREEN)
multi(s, Inches(7.34), Inches(3.02), Inches(4.65), Inches(1.8), [
    ("Skill 固定输入/输出契约", {"size": 18}),
    ("diagnostic_summary 先算事实", {"size": 18, "color": GREEN}),
    ("LLM 数字校验 + fallback", {"size": 18, "color": GREEN}),
    ("Dashboard 展示状态和产物", {"size": 18, "color": GREEN}),
], line_spacing=1.28)

rect(s, Inches(1.20), Inches(6.12), Inches(10.9), Inches(0.78), fill=RGBColor(0xEE, 0xF7, 0xEE), line=GREEN, radius=True)
text(s, Inches(1.55), Inches(6.32), Inches(10.2), Inches(0.32),
     "OpenClaw 把 ML + LLM 组织成可审计、可复现的 agent workflow。", size=18, bold=True, color=INK)
footer(s, "OpenClaw")


# Slide 9: Today Radar
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "Today Radar：把模型用于今天的候选发现", 9)
text(s, Inches(0.72), Inches(1.12), Inches(11.4), Inches(0.42),
     "扫描创建于 30-45 天前的仓库，复用训练 schema 和 RF 模型，输出候选观察清单。", size=21, bold=True)

step(s, Inches(0.86), Inches(2.02), "01", "近期仓库", "created 30-45 days ago", GREEN)
step(s, Inches(0.86), Inches(3.04), "02", "特征对齐", "same 38-feature schema", BLUE)
step(s, Inches(0.86), Inches(4.06), "03", "模型打分", "RF attention_score", ORANGE)
step(s, Inches(0.86), Inches(5.08), "04", "行动分层", "deep_dive / try / watch", PURPLE)

rect(s, Inches(6.50), Inches(2.02), Inches(5.75), Inches(3.60), fill=SUBTLE, line=BORDER, radius=True)
text(s, Inches(6.88), Inches(2.34), Inches(4.9), Inches(0.34), "当前 Top Candidates", size=22, bold=True)
rows = [
    ("#1", "run-llama/ParseBench", "0.795", "deep_dive", GREEN),
    ("#2", "SeanFDZ/macmind", "0.355", "try", BLUE),
    ("#3", "gryszzz/OpenThymos", "0.315", "try", BLUE),
]
for i, (rank, repo, score, action, c) in enumerate(rows):
    y = Inches(3.02 + i * 0.62)
    text(s, Inches(6.88), y, Inches(0.55), Inches(0.26), rank, size=14, bold=True, font=FONT)
    text(s, Inches(7.48), y, Inches(2.65), Inches(0.26), repo, size=14)
    text(s, Inches(10.12), y, Inches(0.75), Inches(0.26), score, size=14, bold=True, font=FONT)
    pill(s, Inches(10.95), y - Inches(0.04), Inches(0.95), action, c)
rect(s, Inches(6.88), Inches(5.05), Inches(4.95), Inches(0.36),
     fill=RGBColor(0xFF, 0xF8, 0xE6), line=ORANGE, radius=True)
text(s, Inches(7.05), Inches(5.135), Inches(4.60), Inches(0.16),
     "候选清单，不是已验证预测", size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
footer(s, "Today Radar")


# Slide 10: Summary
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "总结：这个系统已经完成了一个闭环", 10)

rect(s, Inches(0.82), Inches(1.55), Inches(11.7), Inches(1.00), fill=DARK, line=None, radius=True)
text(s, Inches(1.10), Inches(1.80), Inches(11.1), Inches(0.38),
     "数据挖掘模型 + OpenClaw 工作流 + 今日候选应用", size=27, bold=True, color=BG, align=PP_ALIGN.CENTER)

cards = [
    ("模型层", "早期 30 天信号可以对潜力仓库排序", "RF AUC 0.878 · P@10 0.90", GREEN),
    ("系统层", "OpenClaw 把复杂任务拆成可复现步骤", "skills · artifacts · dashboard", BLUE),
    ("应用层", "Today Radar 给出今天值得看的候选项目", "30-45 天窗口 · 动态阈值", ORANGE),
]
for i, (title, claim, proof, c) in enumerate(cards):
    x = Inches(0.82 + i * 4.05)
    rect(s, x, Inches(3.05), Inches(3.55), Inches(2.30), fill=SUBTLE, line=BORDER, radius=True)
    rect(s, x, Inches(3.05), Inches(3.55), Inches(0.08), fill=c, line=None)
    text(s, x + Inches(0.25), Inches(3.34), Inches(3.05), Inches(0.34), title, size=22, bold=True)
    text(s, x + Inches(0.25), Inches(3.92), Inches(3.05), Inches(0.60), claim, size=16, bold=True)
    text(s, x + Inches(0.25), Inches(4.78), Inches(3.05), Inches(0.24), proof, size=11, color=MUTED, font=FONT)

text(s, Inches(0.82), Inches(6.28), Inches(11.7), Inches(0.40),
     "展示主线：先证明模型有效，再说明 OpenClaw 如何让这件事从一次实验变成可运行的系统。", size=19, bold=True, align=PP_ALIGN.CENTER)
footer(s, "Summary")


prs.save(OUT)
print("Saved:", OUT)
